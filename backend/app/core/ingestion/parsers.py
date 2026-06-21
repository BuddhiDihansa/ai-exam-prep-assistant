"""
Turns an uploaded file into a list of {"text": str, "page_number": int} blocks.
page_number lets us cite "page 4" back to the user later.
"""
import fitz  # PyMuPDF
import docx
from pptx import Presentation


def parse_pdf(path: str) -> list[dict]:
    blocks = []
    with fitz.open(path) as pdf:
        for page_index, page in enumerate(pdf):
            text = page.get_text().strip()
            if text:
                blocks.append({"text": text, "page_number": page_index + 1})
    return blocks


def parse_docx(path: str) -> list[dict]:
    document = docx.Document(path)
    # DOCX has no real "pages" until rendered, so we treat the whole
    # document as page 1 and rely on chunk_index for ordering.
    full_text = "\n".join(p.text for p in document.paragraphs if p.text.strip())
    return [{"text": full_text, "page_number": 1}] if full_text else []


def parse_pptx(path: str) -> list[dict]:
    prs = Presentation(path)
    blocks = []
    for slide_index, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t)
        slide_text = "\n".join(texts).strip()
        if slide_text:
            blocks.append({"text": slide_text, "page_number": slide_index + 1})
    return blocks


PARSERS = {
    "pdf": parse_pdf,
    "docx": parse_docx,
    "pptx": parse_pptx,
}


def parse_file(path: str, file_type: str) -> list[dict]:
    file_type = file_type.lower().lstrip(".")
    if file_type not in PARSERS:
        raise ValueError(f"Unsupported file type: {file_type}")
    return PARSERS[file_type](path)